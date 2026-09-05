"""
AI Report Summary Generator
Generates comprehensive, well-formatted reports with AI insights
"""
import json
import os
from io import BytesIO
from datetime import datetime, timedelta, timezone
from flask import Blueprint, request, jsonify, send_file
from flask_jwt_extended import jwt_required, get_jwt_identity
import logging

from .models import execute_safe_query, get_db
from .ai.trends import analyze_theme_frequency, analyze_sentiment_trends, analyze_theme_trends, detect_theme_anomalies
from .ai.themes import get_theme_trends_by_year
from .ai.themes import _normalize_theme_name

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except Exception:
    Presentation = None
    Inches = None
    Pt = None
    CategoryChartData = None
    XL_CHART_TYPE = None
    XL_LEGEND_POSITION = None
    MSO_SHAPE = None
    PP_ALIGN = None
    MSO_ANCHOR = None
    RGBColor = None

try:
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    from matplotlib import rcParams
    from PIL import Image
    MATPLOTLIB_AVAILABLE = True
except Exception:
    plt = None
    Image = None
    MATPLOTLIB_AVAILABLE = False

report_bp = Blueprint('reports', __name__)
logger = logging.getLogger(__name__)
DEFAULT_PRESENTATION_FONT = 'Plus Jakarta Sans'

PRESENTATION_THEME_STYLES = {
    'corporate': {
        'label': 'Corporate Blue',
        'bg': 'F8FAFF',
        'title': '1E3A8A',
        'text': '0F172A',
        'accent': '2563EB',
        'muted': '475569',
        'positive': '16A34A',
        'neutral': 'F59E0B',
        'negative': 'DC2626',
    },
    'ocean': {
        'label': 'Ocean Teal',
        'bg': 'F2FBFA',
        'title': '0F766E',
        'text': '0B3B39',
        'accent': '14B8A6',
        'muted': '155E75',
        'positive': '16A34A',
        'neutral': 'EA580C',
        'negative': 'B91C1C',
    },
    'sunrise': {
        'label': 'Sunrise Amber',
        'bg': 'FFF8F1',
        'title': '9A3412',
        'text': '431407',
        'accent': 'F97316',
        'muted': '7C2D12',
        'positive': '15803D',
        'neutral': 'CA8A04',
        'negative': 'B91C1C',
    }
}


def _hex_to_rgb(color_hex):
    color = str(color_hex or '').strip().lstrip('#')
    if len(color) != 6:
        color = '334155'
    return RGBColor(int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16))


def _apply_slide_background(slide, color_hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)


def _add_footer(slide, style, text):
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
    tf = footer.text_frame
    tf.text = text
    _set_text_frame_style(tf, size=9, color=style['muted'], italic=True)


def _set_text_frame_style(text_frame, size=None, color=None, bold=None, italic=None, align=None, font_name=DEFAULT_PRESENTATION_FONT):
    for paragraph in text_frame.paragraphs:
        if align is not None and PP_ALIGN is not None:
            paragraph.alignment = align
        if size is not None:
            paragraph.font.size = Pt(size)
        if color is not None:
            paragraph.font.color.rgb = _hex_to_rgb(color)
        if bold is not None:
            paragraph.font.bold = bold
        if italic is not None:
            paragraph.font.italic = italic
        paragraph.font.name = font_name
        for run in paragraph.runs:
            if size is not None:
                run.font.size = Pt(size)
            if color is not None:
                run.font.color.rgb = _hex_to_rgb(color)
            if bold is not None:
                run.font.bold = bold
            if italic is not None:
                run.font.italic = italic
            run.font.name = font_name


def _add_card(slide, left, top, width, height, fill='FFFFFF', line='D8E2F0'):
    if not MSO_SHAPE:
        return slide.shapes.add_textbox(left, top, width, height)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = _hex_to_rgb(fill)
    card.line.color.rgb = _hex_to_rgb(line)
    return card


def _add_top_accent(slide, style):
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _hex_to_rgb(style['accent'])
    accent.line.fill.background()


def _safe_int(value, default=0):
    try:
        return int(value)
    except Exception:
        return default


def _as_bool(value, default=False):
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value != 0
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in {'1', 'true', 'yes', 'y', 'on'}:
            return True
        if normalized in {'0', 'false', 'no', 'n', 'off'}:
            return False
    return default


def _has_nonzero(values):
    return any(_safe_int(v, 0) > 0 for v in values)


def _build_summary_payload(year, limit):
    top_themes = _get_top_themes(year=year, limit=limit)
    sentiment_data = _get_sentiment_analysis(year=year)
    anomalies = _get_anomalies(year=year, limit=limit)
    recommendations = _generate_recommendations(top_themes, sentiment_data, anomalies)

    executive_summary = ""
    if sentiment_data and top_themes:
        top_theme = _display_theme_name(top_themes[0])
        positive_rate = sentiment_data.get('positive_rate', 0)
        executive_summary = (
            f"In {year}, '{top_theme}' was the dominant governance theme with {positive_rate}% positive sentiment. "
            f"Overall trend is {sentiment_data.get('trend', 'stable')}."
        )
    elif top_themes:
        top_theme = _display_theme_name(top_themes[0])
        executive_summary = f"In {year}, '{top_theme}' was the primary governance focus area."

    total_meetings = 0
    try:
        result = execute_safe_query(
            '''
            SELECT COUNT(DISTINCT meeting_id) as count
            FROM Meetings
            WHERE strftime('%Y', meeting_date) = ?
            ''',
            (str(year),)
        )
        total_meetings = result[0].get('count', 0) if result else 0
    except Exception as e:
        logger.warning(f"Could not count meetings: {e}")

    total_mentions_all = sum(theme.get('total_mentions', 0) for theme in top_themes) or 1

    return {
        'year': year,
        'executive_summary': executive_summary,
        'statistics': {
            'total_meetings': total_meetings,
            'total_themes': len(top_themes) if top_themes else 0,
            'sentiment_trend': sentiment_data.get('trend', 'unknown') if sentiment_data else 'unknown',
            'critical_anomalies': len([a for a in anomalies if str(a.get('severity', '')).lower() == 'critical'])
        },
        'top_themes': [
            {
                'theme': _display_theme_name(t),
                'mentions': t.get('total_mentions', 0),
                'percentage': round((t.get('total_mentions', 0) / total_mentions_all) * 100, 1)
            }
            for t in top_themes
        ],
        'sentiment': sentiment_data,
        'critical_anomalies': [
            {
                'theme': _display_theme_name(a),
                'month': a.get('month', 'Unknown'),
                'mentions': a.get('mention_count', 0),
                'baseline': a.get('expected_baseline', 0),
                'z_score': round(float(a.get('z_score', 0)), 2),
                'severity': a.get('severity', 'unknown')
            }
            for a in anomalies
        ],
        'recommendations': recommendations,
        'generated_at': datetime.now(timezone.utc).isoformat()
    }


def _build_presentation_bytes(
    data,
    template_theme='corporate',
    include_anomalies=True,
    include_speaker_notes=True,
    title=None,
    slide_mode='auto',
    include_appendix=False,
    options=None,
):
    if not Presentation:
        raise RuntimeError('python-pptx dependency is not available')

    style = PRESENTATION_THEME_STYLES.get(template_theme, PRESENTATION_THEME_STYLES['corporate'])
    year = data.get('year', datetime.now().year)
    generated_at = datetime.now().strftime('%B %d, %Y %I:%M %p UTC')

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Add baseline metadata for auditability and enterprise records.
    core = prs.core_properties
    core.author = 'ITDS Automated Reporting'
    core.last_modified_by = 'ITDS Service'
    core.title = title or f'Governance Insights Presentation ({year})'
    core.subject = f'Governance analytics summary for {year}'
    core.keywords = 'governance, analytics, sentiment, themes, anomalies'

    normalized_slide_mode = str(slide_mode or 'auto').strip().lower()
    if normalized_slide_mode not in {'auto', 'fixed'}:
        normalized_slide_mode = 'auto'

    theme_rows = data.get('top_themes', [])[:8]
    sentiment = data.get('sentiment') or {}
    anomalies = data.get('critical_anomalies', []) if include_anomalies else []
    recommendations = data.get('recommendations') or []
    options = options or {}

    has_theme_data = len(theme_rows) > 0
    has_sentiment_data = _has_nonzero([
        sentiment.get('positive_rate', 0),
        sentiment.get('neutral_rate', 0),
        sentiment.get('negative_rate', 0),
    ])
    has_anomaly_data = len(anomalies) > 0
    has_recommendation_data = len(recommendations) > 0

    # Slide 1: Title
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    _apply_slide_background(s1, style['bg'])
    _add_top_accent(s1, style)
    s1.shapes.title.text = title or f"Governance Insights — Executive Summary ({year})"
    _set_text_frame_style(s1.shapes.title.text_frame, size=34, color=style['title'], bold=True, align=PP_ALIGN.CENTER)
    subtitle_tf = s1.placeholders[1].text_frame
    subtitle_tf.text = f"Prepared for: Board of Directors • Template: {style['label']} • Generated: {generated_at} • Confidential"
    _set_text_frame_style(subtitle_tf, size=12, color=style['muted'], align=PP_ALIGN.CENTER)
    if include_speaker_notes:
        s1.notes_slide.notes_text_frame.text = 'Present to Board — key findings, implications, and recommended actions.'

    # Slide 2: Executive Summary
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_slide_background(s2, style['bg'])
    _add_top_accent(s2, style)
    s2.shapes.title.text = 'Executive Summary'
    _set_text_frame_style(s2.shapes.title.text_frame, size=28, color=style['title'], bold=True)

    summary_card = _add_card(s2, Inches(0.6), Inches(1.05), Inches(12.1), Inches(1.05), fill='FFFFFF', line='D9E6F7')
    summary_tf = summary_card.text_frame
    summary_tf.margin_left = Inches(0.25)
    summary_tf.margin_right = Inches(0.25)
    summary_tf.margin_top = Inches(0.12)
    summary_tf.word_wrap = True
    summary_tf.text = data.get('executive_summary', 'No summary available.')
    _set_text_frame_style(summary_tf, size=16, color=style['text'])

    insight_card = _add_card(s2, Inches(0.6), Inches(2.2), Inches(12.1), Inches(0.55), fill='EEF4FF', line='D9E6F7')
    insight_tf = insight_card.text_frame
    insight_tf.text = f"Year focus: {year} • Template: {style['label']} • Output optimized for presentation review"
    insight_tf.margin_left = Inches(0.25)
    insight_tf.margin_top = Inches(0.08)
    _set_text_frame_style(insight_tf, size=11, color=style['muted'], italic=True)

    # Stats grid - 4 columns
    stats = data.get('statistics', {})
    stat_data = [
        ('Total Meetings', stats.get('total_meetings', 0)),
        ('Unique Themes', stats.get('total_themes', 0)),
        ('Sentiment Trend', stats.get('sentiment_trend', 'unknown').capitalize()),
        ('Critical Alerts', stats.get('critical_anomalies', 0)),
    ]
    
    col_width = Inches(2.9)
    start_y = Inches(3.05)
    start_x = Inches(0.6)
    
    for idx, (label, value) in enumerate(stat_data):
        col_x = start_x + (idx * col_width)
        card = _add_card(s2, col_x, start_y, col_width - Inches(0.15), Inches(0.95), fill='FFFFFF', line='D9E6F7')
        tf = card.text_frame
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.1)
        tf.word_wrap = True
        tf.text = label
        _set_text_frame_style(tf, size=10, color=style['muted'], bold=True)
        value_para = tf.add_paragraph()
        value_para.text = str(value)
        value_para.space_before = Pt(4)
        value_para.font.size = Pt(24)
        value_para.font.bold = True
        value_para.font.color.rgb = _hex_to_rgb(style['accent'])
        value_para.font.name = DEFAULT_PRESENTATION_FONT
    
    _add_footer(s2, style, f"Data source: ITDS analytics APIs | Year: {year}")
    if include_speaker_notes:
        s2.notes_slide.notes_text_frame.text = 'Summarize what changed this year and what needs attention next.'

    if normalized_slide_mode == 'fixed' or has_theme_data:
        # Slide: Top Themes chart (rendered as image for higher fidelity)
        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s3, style['bg'])
        _add_top_accent(s3, style)
        s3.shapes.title.text = 'Top Themes by Mentions'
        _set_text_frame_style(s3.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        try:
            if MATPLOTLIB_AVAILABLE and theme_rows:
                labels = [row.get('theme', 'Unknown')[:40] for row in theme_rows]
                values = [max(0, _safe_int(row.get('mentions'))) for row in theme_rows]
                fig, ax = plt.subplots(figsize=(8.8, 4.9))
                ordered_labels = labels[::-1]
                ordered_values = values[::-1]
                bars = ax.barh(ordered_labels, ordered_values, color=f"#{style['accent']}", height=0.58)
                ax.set_xlabel('Mentions')
                ax.set_frame_on(False)
                ax.tick_params(axis='x', labelsize=9, colors='#475569')
                ax.tick_params(axis='y', labelsize=9, colors='#0f172a')
                ax.xaxis.grid(True, linestyle='--', linewidth=0.7, alpha=0.22)
                ax.set_axisbelow(True)
                ax.spines['top'].set_visible(False)
                ax.spines['right'].set_visible(False)
                ax.spines['left'].set_visible(False)
                ax.spines['bottom'].set_color('#cbd5e1')
                max_value = max(ordered_values) if ordered_values else 0
                for bar, value in zip(bars, ordered_values):
                    ax.text(value + (max_value * 0.02 if max_value else 0.1), bar.get_y() + bar.get_height() / 2, str(value), va='center', ha='left', fontsize=9, color='#1e3a8a', fontweight='bold')
                plt.tight_layout()
                img_buf = BytesIO()
                fig.savefig(img_buf, format='png', dpi=150, bbox_inches='tight', pad_inches=0.2)
                plt.close(fig)
                img_buf.seek(0)
                chart_card = _add_card(s3, Inches(0.6), Inches(1.05), Inches(8.2), Inches(5.55), fill='FFFFFF', line='D9E6F7')
                s3.shapes.add_picture(img_buf, Inches(0.85), Inches(1.25), Inches(7.75), Inches(5.0))
            else:
                # Fallback: small text summary if matplotlib not available
                chart_card = _add_card(s3, Inches(0.6), Inches(1.05), Inches(8.2), Inches(5.55), fill='FFFFFF', line='D9E6F7')
                summary_tf = chart_card.text_frame
                summary_tf.margin_left = Inches(0.18)
                summary_tf.margin_top = Inches(0.14)
                summary_tf.word_wrap = True
                summary_tf.text = '\n'.join([f"{i+1}. {r.get('theme', 'Unknown')} — {r.get('mentions', 0)}" for i, r in enumerate(theme_rows[:8])])
                _set_text_frame_style(summary_tf, size=12, color=style['text'])
        except Exception:
            logger.exception('Could not render top themes chart as image, falling back to text')

        side_card = _add_card(s3, Inches(9.0), Inches(1.05), Inches(3.0), Inches(5.55), fill='F8FBFF', line='D9E6F7')
        side_tf = side_card.text_frame
        side_tf.margin_left = Inches(0.18)
        side_tf.margin_top = Inches(0.16)
        side_tf.word_wrap = True
        side_tf.text = 'Why it matters'
        _set_text_frame_style(side_tf, size=13, color=style['title'], bold=True)
        for idx, row in enumerate(theme_rows[:4], 1):
            p = side_tf.add_paragraph()
            p.text = f"{idx}. {row.get('theme', 'Unknown')}"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(style['text'])
            p.font.name = DEFAULT_PRESENTATION_FONT
            detail = side_tf.add_paragraph()
            detail.text = f"{row.get('mentions', 0)} mentions"
            detail.font.size = Pt(10)
            detail.font.color.rgb = _hex_to_rgb(style['muted'])
            detail.font.name = DEFAULT_PRESENTATION_FONT

        _add_footer(s3, style, f"Top themes limited to {len(theme_rows)} items")
        if include_speaker_notes:
            s3.notes_slide.notes_text_frame.text = 'Highlight top themes and explain why mention count changed.'

    if normalized_slide_mode == 'fixed' or has_sentiment_data:
        # Slide: Sentiment breakdown chart (render as image for crisp visuals)
        s4 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s4, style['bg'])
        _add_top_accent(s4, style)
        s4.shapes.title.text = 'Sentiment Breakdown'
        _set_text_frame_style(s4.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        try:
            if MATPLOTLIB_AVAILABLE:
                labels = ['Positive', 'Neutral', 'Negative']
                values = [max(0, _safe_int(sentiment.get('positive_rate', 0))), max(0, _safe_int(sentiment.get('neutral_rate', 0))), max(0, _safe_int(sentiment.get('negative_rate', 0)))]
                fig, ax = plt.subplots(figsize=(5.0, 4.9))
                wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.0f%%', startangle=90, colors=[f"#{style['positive']}", f"#{style['neutral']}", f"#{style['negative']}"], textprops={'fontsize': 10})
                ax.axis('equal')
                centre_circle = plt.Circle((0, 0), 0.62, fc='white')
                fig.gca().add_artist(centre_circle)
                plt.tight_layout()
                img_buf = BytesIO()
                fig.savefig(img_buf, format='png', dpi=150, bbox_inches='tight', pad_inches=0.2)
                plt.close(fig)
                img_buf.seek(0)
                chart_card = _add_card(s4, Inches(0.6), Inches(1.05), Inches(5.15), Inches(5.55), fill='FFFFFF', line='D9E6F7')
                s4.shapes.add_picture(img_buf, Inches(0.8), Inches(1.28), Inches(4.75), Inches(4.9))
            else:
                # Fallback: textual summary
                chart_card = _add_card(s4, Inches(0.6), Inches(1.05), Inches(5.15), Inches(5.55), fill='FFFFFF', line='D9E6F7')
                details_box = chart_card.text_frame
                details_box.word_wrap = True
                details_box.text = f"Positive: {sentiment.get('positive_rate', 0)}%\nNeutral: {sentiment.get('neutral_rate', 0)}%\nNegative: {sentiment.get('negative_rate', 0)}%"
                _set_text_frame_style(details_box, size=14, color=style['text'])
        except Exception:
            logger.exception('Could not render sentiment chart as image, falling back to text')

        details_card = _add_card(s4, Inches(6.0), Inches(1.05), Inches(6.3), Inches(5.55), fill='F8FBFF', line='D9E6F7')
        details_box = details_card.text_frame
        details_box.word_wrap = True
        details_box.margin_left = Inches(0.2)
        details_box.margin_top = Inches(0.16)

        trend_para = details_box.paragraphs[0]
        trend_para.text = 'Overall Trend'
        trend_para.font.size = Pt(13)
        trend_para.font.color.rgb = _hex_to_rgb(style['muted'])
        trend_para.font.bold = True
        trend_para.font.name = DEFAULT_PRESENTATION_FONT

        trend_value = details_box.add_paragraph()
        trend_value.text = sentiment.get('trend', 'unknown').capitalize()
        trend_value.font.size = Pt(26)
        trend_value.font.bold = True
        trend_value.font.color.rgb = _hex_to_rgb(style['accent'])
        trend_value.space_before = Pt(4)
        trend_value.space_after = Pt(16)
        trend_value.font.name = DEFAULT_PRESENTATION_FONT

        interpretation = details_box.add_paragraph()
        interpretation.text = 'The balance is heavily neutral, with a positive tilt and no visible negative pressure.'
        interpretation.font.size = Pt(12)
        interpretation.font.color.rgb = _hex_to_rgb(style['text'])
        interpretation.space_after = Pt(10)
        interpretation.font.name = DEFAULT_PRESENTATION_FONT

        for label, key, color in [
            ('Positive', 'positive_rate', style['positive']),
            ('Neutral', 'neutral_rate', style['neutral']),
            ('Negative', 'negative_rate', style['negative']),
        ]:
            p = details_box.add_paragraph()
            p.text = f"{label}: {sentiment.get(key, 0)}%"
            p.font.size = Pt(14)
            p.font.color.rgb = _hex_to_rgb(color)
            p.space_before = Pt(3)
            p.space_after = Pt(3)
            p.font.name = DEFAULT_PRESENTATION_FONT

        _add_footer(s4, style, 'Sentiment source: /api/ai/sentiment-trends')
        if include_speaker_notes:
            s4.notes_slide.notes_text_frame.text = 'Connect sentiment changes to major themes and events.'

    if normalized_slide_mode == 'fixed' or has_anomaly_data or has_recommendation_data:
        # Slide: Anomalies and recommendations
        s5 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s5, style['bg'])
        _add_top_accent(s5, style)
        s5.shapes.title.text = 'Anomalies and Recommendations'
        _set_text_frame_style(s5.shapes.title.text_frame, size=28, color=style['title'], bold=True)
        
        # LEFT COLUMN: Anomalies
        anom_card = _add_card(s5, Inches(0.6), Inches(1.15), Inches(5.8), Inches(5.3), fill='FFFFFF', line='D9E6F7')
        anom_box = anom_card.text_frame
        anom_box.word_wrap = True
        anom_box.margin_left = Inches(0.18)
        anom_box.margin_top = Inches(0.12)
        anom_box.text = 'Critical Anomalies'
        _set_text_frame_style(anom_box, size=13, color=style['accent'], bold=True)
        
        if anomalies:
            for item in anomalies[:3]:
                theme_name = item.get('theme', 'Unknown')
                p = anom_box.add_paragraph()
                p.text = f"• {theme_name}"
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = _hex_to_rgb(style['text'])
                p.space_before = Pt(4)
                p.space_after = Pt(2)
                p.font.name = DEFAULT_PRESENTATION_FONT
                
                detail = anom_box.add_paragraph()
                detail.text = f"{item.get('month', 'N/A')} • {item.get('mentions', 0)} mentions (z={item.get('z_score', 0)})"
                detail.font.size = Pt(10)
                detail.font.color.rgb = _hex_to_rgb(style['muted'])
                detail.level = 1
                detail.space_after = Pt(6)
                detail.font.name = DEFAULT_PRESENTATION_FONT
        else:
            p = anom_box.add_paragraph()
            p.text = 'No critical anomalies detected.'
            p.font.size = Pt(11)
            p.font.color.rgb = _hex_to_rgb(style['muted'])
            p.font.name = DEFAULT_PRESENTATION_FONT
        
        # RIGHT COLUMN: Recommendations
        rec_card = _add_card(s5, Inches(6.8), Inches(1.15), Inches(5.9), Inches(5.3), fill='F8FBFF', line='D9E6F7')
        rec_box = rec_card.text_frame
        rec_box.word_wrap = True
        rec_box.margin_left = Inches(0.18)
        rec_box.margin_top = Inches(0.12)
        rec_box.text = 'Key Actions'
        _set_text_frame_style(rec_box, size=13, color=style['accent'], bold=True)
        
        if recommendations:
            for idx, rec in enumerate(recommendations[:3], 1):
                p = rec_box.add_paragraph()
                p.text = f"{idx}. {rec}"
                p.font.size = Pt(11)
                p.font.color.rgb = _hex_to_rgb(style['text'])
                p.space_before = Pt(3)
                p.space_after = Pt(6)
                p.font.name = DEFAULT_PRESENTATION_FONT
        else:
            p = rec_box.add_paragraph()
            p.text = 'No recommendations generated.'
            p.font.size = Pt(11)
            p.font.color.rgb = _hex_to_rgb(style['muted'])
            p.font.name = DEFAULT_PRESENTATION_FONT
        
        _add_footer(s5, style, 'Generated recommendations should be reviewed by meeting owners')
        if include_speaker_notes:
            s5.notes_slide.notes_text_frame.text = 'Walk through risks first, then focus on top 2 actionable recommendations.'

    if include_appendix:
        # Optional appendix slide for governance/audit metadata.
        s7 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s7, style['bg'])
        _add_top_accent(s7, style)
        s7.shapes.title.text = 'Appendix: Data Quality and Scope'
        _set_text_frame_style(s7.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        appendix_card = _add_card(s7, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.55), fill='FFFFFF', line='D9E6F7')
        appendix = appendix_card.text_frame
        appendix.word_wrap = True
        appendix.margin_left = Inches(0.22)
        appendix.margin_top = Inches(0.18)
        appendix.text = 'Report metadata'
        _set_text_frame_style(appendix, size=13, color=style['accent'], bold=True)
        
        appendix_metadata = [
            ('Reporting Year', str(year)),
            ('Themes Analyzed', str(len(theme_rows))),
            ('Sentiment Data', 'Yes' if has_sentiment_data else 'No'),
            ('Anomaly Detection', 'Yes' if include_anomalies else 'No'),
            ('Slide Strategy', normalized_slide_mode.capitalize()),
            ('Generated', generated_at),
        ]
        
        for idx, (label, value) in enumerate(appendix_metadata):
            p = appendix.add_paragraph()
            p.text = f"{label}: "
            p.font.size = Pt(12)
            p.font.color.rgb = _hex_to_rgb(style['muted'])
            p.font.name = DEFAULT_PRESENTATION_FONT

            run = p.add_run()
            run.text = value
            run.font.bold = True
            run.font.color.rgb = _hex_to_rgb(style['accent'])
            run.font.name = DEFAULT_PRESENTATION_FONT

            p.space_before = Pt(4)
            p.space_after = Pt(4)
        
        _add_footer(s7, style, 'Appendix for review and compliance documentation')
        if include_speaker_notes:
            s7.notes_slide.notes_text_frame.text = 'Use this appendix when sharing outside the analytics team.'

    if options.get('include_sentiment_trends'):
        s8 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s8, style['bg'])
        _add_top_accent(s8, style)
        s8.shapes.title.text = 'Sentiment Trend Overview'
        _set_text_frame_style(s8.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        summary_card = _add_card(s8, Inches(0.6), Inches(1.15), Inches(3.3), Inches(5.3), fill='FFFFFF', line='D9E6F7')
        summary_tf = summary_card.text_frame
        summary_tf.word_wrap = True
        summary_tf.margin_left = Inches(0.18)
        summary_tf.margin_top = Inches(0.15)
        summary_tf.text = 'Snapshot'
        _set_text_frame_style(summary_tf, size=13, color=style['accent'], bold=True)
        for label, value, color in [
            ('Positive', sentiment.get('positive_rate', 0), style['positive']),
            ('Neutral', sentiment.get('neutral_rate', 0), style['neutral']),
            ('Negative', sentiment.get('negative_rate', 0), style['negative']),
        ]:
            p = summary_tf.add_paragraph()
            p.text = f'{label}: {value}%'
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(color)
            p.font.name = DEFAULT_PRESENTATION_FONT
            p.space_after = Pt(8)

        insight_card = _add_card(s8, Inches(4.1), Inches(1.15), Inches(8.6), Inches(5.3), fill='F8FBFF', line='D9E6F7')
        insight_tf = insight_card.text_frame
        insight_tf.word_wrap = True
        insight_tf.margin_left = Inches(0.2)
        insight_tf.margin_top = Inches(0.15)
        insight_tf.text = 'Interpretation'
        _set_text_frame_style(insight_tf, size=13, color=style['title'], bold=True)

        trend = str(sentiment.get('trend', 'stable')).capitalize()
        p = insight_tf.add_paragraph()
        p.text = f'Trend: {trend}'
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = _hex_to_rgb(style['accent'])
        p.font.name = DEFAULT_PRESENTATION_FONT
        p.space_after = Pt(10)

        p = insight_tf.add_paragraph()
        p.text = 'The sentiment profile is mostly neutral, with a stable-to-improving tone.'
        p.font.size = Pt(14)
        p.font.color.rgb = _hex_to_rgb(style['text'])
        p.font.name = DEFAULT_PRESENTATION_FONT
        p.space_after = Pt(4)
        _add_footer(s8, style, f'Sentiment snapshot for {year}')
        if include_speaker_notes:
            s8.notes_slide.notes_text_frame.text = 'Use this slide to summarize overall sentiment direction.'

    if options.get('include_growth_analysis'):
        s9 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s9, style['bg'])
        _add_top_accent(s9, style)
        s9.shapes.title.text = 'Growth Analysis'
        _set_text_frame_style(s9.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        growth_card = _add_card(s9, Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.3), fill='FFFFFF', line='D9E6F7')
        growth_box = growth_card.text_frame
        growth_box.word_wrap = True
        growth_box.margin_left = Inches(0.18)
        growth_box.margin_top = Inches(0.15)
        growth_box.text = 'Year-over-year movement'
        _set_text_frame_style(growth_box, size=13, color=style['accent'], bold=True)

        growth_items = options.get('growth_data') or []
        for idx, item in enumerate(growth_items[:5], 1):
            p = growth_box.add_paragraph()
            p.text = f"{idx}. {item.get('name') or item.get('theme_id', 'Theme')}"
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(style['text'])
            p.font.name = DEFAULT_PRESENTATION_FONT
            p.space_after = Pt(1)

            detail = growth_box.add_paragraph()
            detail.text = f"Current: {item.get('current', 0)}  |  Previous: {item.get('previous', 0)}  |  Growth: {item.get('growth_pct', 0)}%  |  Trend: {item.get('trend', '→')}"
            detail.font.size = Pt(11)
            detail.font.color.rgb = _hex_to_rgb(style['muted'])
            detail.font.name = DEFAULT_PRESENTATION_FONT
            detail.space_after = Pt(6)

        if not growth_items:
            p = growth_box.add_paragraph()
            p.text = 'No growth data available for the selected period.'
            p.font.size = Pt(15)
            p.font.color.rgb = _hex_to_rgb(style['muted'])
            p.font.name = DEFAULT_PRESENTATION_FONT
        _add_footer(s9, style, f'Year-over-year theme comparison for {year}')

    if options.get('include_anomaly_details'):
        s10 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s10, style['bg'])
        _add_top_accent(s10, style)
        s10.shapes.title.text = 'Detailed Anomalies'
        _set_text_frame_style(s10.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        anomaly_card = _add_card(s10, Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.3), fill='FFFFFF', line='D9E6F7')
        anomaly_box = anomaly_card.text_frame
        anomaly_box.word_wrap = True
        anomaly_box.margin_left = Inches(0.18)
        anomaly_box.margin_top = Inches(0.15)
        anomaly_box.text = 'Top flagged items'
        _set_text_frame_style(anomaly_box, size=13, color=style['accent'], bold=True)

        anomaly_items = (options.get('anomaly_data') or [])[:6]
        for item in anomaly_items:
            p = anomaly_box.add_paragraph()
            p.text = f"{item.get('theme_id', 'Theme')}"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(style['text'])
            p.font.name = DEFAULT_PRESENTATION_FONT
            p.space_after = Pt(1)

            detail = anomaly_box.add_paragraph()
            detail.text = f"{item.get('month', 'Unknown')}  |  Mentions: {item.get('mentions', 0)}  |  Z: {item.get('z_score', 0)}  |  Severity: {item.get('severity', 'Low')}"
            detail.font.size = Pt(11)
            detail.font.color.rgb = _hex_to_rgb(style['muted'])
            detail.font.name = DEFAULT_PRESENTATION_FONT
            detail.space_after = Pt(6)

        if not anomaly_items:
            p = anomaly_box.add_paragraph()
            p.text = 'No anomaly details were found for this year.'
            p.font.size = Pt(15)
            p.font.color.rgb = _hex_to_rgb(style['muted'])
            p.font.name = DEFAULT_PRESENTATION_FONT
        _add_footer(s10, style, f'Anomaly breakdown for {year}')

    if options.get('include_prioritized_recs') and has_recommendation_data:
        s11 = prs.slides.add_slide(prs.slide_layouts[5])
        _apply_slide_background(s11, style['bg'])
        _add_top_accent(s11, style)
        s11.shapes.title.text = 'Prioritized Recommendations'
        _set_text_frame_style(s11.shapes.title.text_frame, size=28, color=style['title'], bold=True)

        rec_card = _add_card(s11, Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.3), fill='F8FBFF', line='D9E6F7')
        rec_box = rec_card.text_frame
        rec_box.word_wrap = True
        rec_box.margin_left = Inches(0.18)
        rec_box.margin_top = Inches(0.15)
        rec_box.text = 'Priority actions'
        _set_text_frame_style(rec_box, size=13, color=style['accent'], bold=True)
        for idx, rec in enumerate(recommendations[:6], 1):
            p = rec_box.add_paragraph()
            p.text = f'{idx}. {rec}'
            p.font.size = Pt(13)
            p.font.color.rgb = _hex_to_rgb(style['text'])
            p.font.name = DEFAULT_PRESENTATION_FONT
            p.space_after = Pt(5)
        _add_footer(s11, style, f'Priority actions for {year}')

    # Slide 6: Next steps
    # Slide 6: Next steps
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(s6, style['bg'])
    _add_top_accent(s6, style)
    
    title_box = s6.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(12.1), Inches(1.0)).text_frame
    title_box.text = 'Next Steps'
    _set_text_frame_style(title_box, size=30, color=style['title'], bold=True)

    step_cards = [
        (Inches(1.15), Inches(2.0), '1', 'Validate high-impact anomalies with meeting owners.'),
        (Inches(1.15), Inches(3.15), '2', 'Assign owners for top themes and define measurable outcomes.'),
        (Inches(1.15), Inches(4.3), '3', 'Review progress in the next reporting cycle.'),
    ]

    for left, top, step_num, step_text in step_cards:
        card = _add_card(s6, left, top, Inches(11.1), Inches(0.9), fill='FFFFFF', line='D9E6F7')
        tf = card.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)
        tf.word_wrap = True
        tf.text = f"Step {step_num}"
        _set_text_frame_style(tf, size=11, color=style['muted'], bold=True)
        body = tf.add_paragraph()
        body.text = step_text
        body.font.size = Pt(15)
        body.font.color.rgb = _hex_to_rgb(style['text'])
        body.font.name = DEFAULT_PRESENTATION_FONT
    
    _add_footer(s6, style, f"Generated by ITDS • {generated_at}")
    if include_speaker_notes:
        s6.notes_slide.notes_text_frame.text = 'Close with ownership, timeline, and measurable outcomes.'

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def _get_sentiment_trends(year):
    try:
        trends = analyze_sentiment_trends(year=year) or []
        formatted = []
        for row in trends:
            formatted.append({
                'month': row.get('month'),
                'positive_rate': row.get('positive_rate', 0),
                'neutral_rate': row.get('neutral', 0),
                'negative_rate': row.get('negative', 0),
                'total': row.get('total', 0),
            })
        return formatted
    except Exception as e:
        logger.warning(f'Could not load sentiment trends: {e}')
        return []


def _get_growth_analysis(year):
    try:
        current_rows = analyze_theme_frequency(year=year, top_n=10) or []
        prev_rows = analyze_theme_frequency(year=year - 1, top_n=10) or []
        prev_map = {}
        for row in prev_rows:
            key = str(row.get('theme_id') or row.get('name') or '').strip().lower()
            if key:
                prev_map[key] = int(row.get('total_mentions') or 0)

        results = []
        for row in current_rows:
            theme_name = str(row.get('name') or row.get('theme_id') or '').strip()
            key = theme_name.lower()
            current_mentions = int(row.get('total_mentions') or 0)
            previous_mentions = prev_map.get(key, 0)
            growth_pct = ((current_mentions - previous_mentions) / previous_mentions * 100) if previous_mentions else (100 if current_mentions else 0)
            results.append({
                'theme_id': row.get('theme_id') or key,
                'name': theme_name,
                'current': current_mentions,
                'previous': previous_mentions,
                'growth_pct': round(growth_pct, 1),
                'trend': '↑' if growth_pct > 10 else ('↓' if growth_pct < -10 else '→'),
            })
        return sorted(results, key=lambda x: x['growth_pct'], reverse=True)
    except Exception as e:
        logger.warning(f'Could not load growth analysis: {e}')
        return []


def _get_anomaly_details(year):
    try:
        anomalies = detect_theme_anomalies(year=year) or []
        formatted = []
        for row in anomalies:
            severity = str(row.get('severity') or 'low').strip().lower()
            formatted.append({
                'theme_id': row.get('theme_id') or row.get('theme') or 'unknown',
                'theme': row.get('theme') or row.get('theme_id') or 'Unknown',
                'month': row.get('month'),
                'mentions': row.get('mention_count', 0),
                'baseline': row.get('expected_baseline', 0),
                'z_score': row.get('z_score', 0),
                'severity': severity.title(),
                'note': row.get('note', ''),
            })
        return formatted
    except Exception as e:
        logger.warning(f'Could not load anomaly details: {e}')
        return []


def _get_db_connection():
    return get_db()


@report_bp.route('/api/reports/presentation/advanced', methods=['POST'])
@jwt_required()
def export_presentation_advanced():
    try:
        payload = request.get_json(silent=True) or {}
        year = _safe_int(payload.get('year'), datetime.now().year)
        template_theme = str(payload.get('template_theme') or 'corporate').strip().lower()
        include_anomalies = _as_bool(payload.get('include_anomalies'), True)
        include_speaker_notes = _as_bool(payload.get('include_speaker_notes'), True)
        slide_mode = str(payload.get('slide_mode') or 'auto').strip().lower()

        data = _build_summary_payload(year=year, limit=max(3, min(10, _safe_int(payload.get('top_n_themes'), 6))))
        options = {
            'include_sentiment_trends': _as_bool(payload.get('include_sentiment_trends'), True),
            'include_growth_analysis': _as_bool(payload.get('include_growth_analysis'), True),
            'include_anomaly_details': _as_bool(payload.get('include_anomaly_details'), True),
            'include_prioritized_recs': _as_bool(payload.get('include_prioritized_recommendations'), True),
            'growth_data': _get_growth_analysis(year),
            'anomaly_data': _get_anomaly_details(year),
        }

        pptx_stream = _build_presentation_bytes(
            data=data,
            template_theme=template_theme,
            include_anomalies=include_anomalies,
            include_speaker_notes=include_speaker_notes,
            slide_mode=slide_mode,
            options=options,
        )

        safe_theme = template_theme if template_theme in PRESENTATION_THEME_STYLES else 'corporate'
        filename = f'governance-presentation-{year}-{safe_theme}.pptx'
        return send_file(
            pptx_stream,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        logger.error(f'Enhanced PPTX export error: {e}')
        return jsonify({'error': 'Failed to generate enhanced presentation'}), 500


@report_bp.route('/api/reports/schedule', methods=['POST'])
@jwt_required()
def create_scheduled_report():
    try:
        user_id = get_jwt_identity()
        payload = request.get_json(silent=True) or {}
        recipients = payload.get('email_recipients') or []
        if isinstance(recipients, str):
            recipients = [r.strip() for r in recipients.split(',') if r.strip()]
        if not recipients:
            return jsonify({'error': 'At least one email recipient is required'}), 400

        conn = _get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            '''
            INSERT INTO ScheduledReports
            (user_id, report_name, email_recipients, frequency, send_time, template_theme, year)
            VALUES (?, ?, ?, ?, ?, ?, ?)
            ''',
            (
                user_id,
                payload.get('report_name') or 'Automated Report',
                ','.join(recipients),
                payload.get('frequency') or 'monthly',
                payload.get('send_time') or '09:00',
                payload.get('template_theme') or 'corporate',
                _safe_int(payload.get('year'), datetime.now().year),
            )
        )
        conn.commit()
        schedule_id = cursor.lastrowid
        conn.close()
        return jsonify({'status': 'scheduled', 'id': schedule_id}), 201
    except Exception as e:
        logger.error(f'Schedule creation error: {e}')
        return jsonify({'error': 'Failed to create scheduled report'}), 500


@report_bp.route('/api/reports/schedule', methods=['GET'])
@jwt_required()
def get_scheduled_reports():
    try:
        user_id = get_jwt_identity()
        conn = _get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            '''
            SELECT id, report_name, email_recipients, frequency, send_time, template_theme, year, is_active, next_send
            FROM ScheduledReports
            WHERE user_id = ?
            ORDER BY created_at DESC
            ''',
            (user_id,)
        )
        rows = cursor.fetchall()
        conn.close()
        return jsonify([
            {
                'id': row['id'],
                'report_name': row['report_name'],
                'email_recipients': [r.strip() for r in (row['email_recipients'] or '').split(',') if r.strip()],
                'frequency': row['frequency'],
                'send_time': row['send_time'],
                'template_theme': row['template_theme'],
                'year': row['year'],
                'is_active': bool(row['is_active']),
                'next_send': row['next_send'],
            } for row in rows
        ]), 200
    except Exception as e:
        logger.error(f'Get schedules error: {e}')
        return jsonify({'error': 'Failed to retrieve scheduled reports'}), 500


@report_bp.route('/api/reports/schedule/<int:schedule_id>', methods=['PUT'])
@jwt_required()
def update_scheduled_report(schedule_id):
    try:
        user_id = get_jwt_identity()
        payload = request.get_json(silent=True) or {}
        conn = _get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT user_id FROM ScheduledReports WHERE id = ?', (schedule_id,))
        row = cursor.fetchone()
        if not row or row['user_id'] != user_id:
            conn.close()
            return jsonify({'error': 'Not found or unauthorized'}), 404

        fields = []
        values = []
        for key in ('report_name', 'frequency', 'send_time', 'template_theme'):
            if key in payload:
                fields.append(f'{key} = ?')
                values.append(payload[key])
        if 'email_recipients' in payload:
            recipients = payload['email_recipients']
            if isinstance(recipients, list):
                recipients = ','.join(recipients)
            fields.append('email_recipients = ?')
            values.append(recipients)
        if 'is_active' in payload:
            fields.append('is_active = ?')
            values.append(1 if _as_bool(payload.get('is_active'), True) else 0)
        if fields:
            fields.append('updated_at = CURRENT_TIMESTAMP')
            values.append(schedule_id)
            cursor.execute(f"UPDATE ScheduledReports SET {', '.join(fields)} WHERE id = ?", values)
            conn.commit()
        conn.close()
        return jsonify({'status': 'updated', 'id': schedule_id}), 200
    except Exception as e:
        logger.error(f'Update schedule error: {e}')
        return jsonify({'error': 'Failed to update schedule'}), 500


@report_bp.route('/api/reports/schedule/<int:schedule_id>', methods=['DELETE'])
@jwt_required()
def delete_scheduled_report(schedule_id):
    try:
        user_id = get_jwt_identity()
        conn = _get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT user_id FROM ScheduledReports WHERE id = ?', (schedule_id,))
        row = cursor.fetchone()
        if not row or row['user_id'] != user_id:
            conn.close()
            return jsonify({'error': 'Not found or unauthorized'}), 404
        cursor.execute('DELETE FROM ScheduledReports WHERE id = ?', (schedule_id,))
        conn.commit()
        conn.close()
        return jsonify({'status': 'deleted', 'id': schedule_id}), 200
    except Exception as e:
        logger.error(f'Delete schedule error: {e}')
        return jsonify({'error': 'Failed to delete schedule'}), 500


@report_bp.route('/api/reports/sentiment-trends/<int:year>', methods=['GET'])
@jwt_required()
def api_sentiment_trends(year):
    return jsonify(_get_sentiment_trends(year)), 200


@report_bp.route('/api/reports/growth-analysis/<int:year>', methods=['GET'])
@jwt_required()
def api_growth_analysis(year):
    return jsonify(_get_growth_analysis(year)), 200


@report_bp.route('/api/reports/anomalies-detailed/<int:year>', methods=['GET'])
@jwt_required()
def api_anomalies_detailed(year):
    return jsonify(_get_anomaly_details(year)), 200


@report_bp.route('/api/reports/branding', methods=['GET'])
@jwt_required()
def get_branding():
    try:
        user_id = get_jwt_identity()
        conn = _get_db_connection()
        cursor = conn.cursor()
        cursor.execute(
            '''
            SELECT organization_name, logo_url, primary_color, secondary_color, accent_color, watermark, footer_text
            FROM PresentationBranding
            WHERE user_id = ?
            ''',
            (user_id,)
        )
        row = cursor.fetchone()
        conn.close()
        if not row:
            return jsonify({
                'organization_name': '', 'logo_url': '', 'primary_color': '#667eea',
                'secondary_color': '#764ba2', 'accent_color': '#f59e0b', 'watermark': '', 'footer_text': ''
            }), 200
        return jsonify(dict(row)), 200
    except Exception as e:
        logger.error(f'Get branding error: {e}')
        return jsonify({'error': 'Failed to retrieve branding'}), 500


@report_bp.route('/api/reports/branding', methods=['PUT'])
@jwt_required()
def update_branding():
    try:
        user_id = get_jwt_identity()
        payload = request.get_json(silent=True) or {}
        conn = _get_db_connection()
        cursor = conn.cursor()
        cursor.execute('SELECT id FROM PresentationBranding WHERE user_id = ?', (user_id,))
        exists = cursor.fetchone() is not None
        values = (
            payload.get('organization_name', ''),
            payload.get('logo_url', ''),
            payload.get('primary_color', '#667eea'),
            payload.get('secondary_color', '#764ba2'),
            payload.get('accent_color', '#f59e0b'),
            payload.get('watermark', ''),
            payload.get('footer_text', ''),
            user_id,
        )
        if exists:
            cursor.execute(
                '''
                UPDATE PresentationBranding
                SET organization_name = ?, logo_url = ?, primary_color = ?, secondary_color = ?, accent_color = ?, watermark = ?, footer_text = ?, updated_at = CURRENT_TIMESTAMP
                WHERE user_id = ?
                ''',
                values
            )
        else:
            cursor.execute(
                '''
                INSERT INTO PresentationBranding
                (organization_name, logo_url, primary_color, secondary_color, accent_color, watermark, footer_text, user_id)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                ''',
                values
            )
        conn.commit()
        conn.close()
        return jsonify({'status': 'updated'}), 200
    except Exception as e:
        logger.error(f'Update branding error: {e}')
        return jsonify({'error': 'Failed to update branding'}), 500


def _get_date_range(days_back=30):
    """Get date range for analysis."""
    end_date = datetime.now(timezone.utc)
    start_date = end_date - timedelta(days=days_back)
    return start_date.isoformat(), end_date.isoformat()


def _get_top_themes(year=None, limit=5):
    """Get top themes for the year or period."""
    try:
        themes = analyze_theme_frequency(year=year)
        if not themes:
            return []
        return themes[:limit]
    except Exception as e:
        logger.warning(f"Could not analyze theme frequency: {e}")
        return []


def _display_theme_name(theme, default='Cluster'):
    name = str(theme.get('name') or theme.get('theme') or '').strip()
    if name and name.lower() not in {'theme', 'themes', 'unknown', 'unknown theme', 'general', 'general topic', 'topic'}:
        return name

    keywords = [str(k).strip() for k in (theme.get('keywords') or []) if str(k).strip()]
    if keywords:
        candidate = ' '.join(keywords[:4]).strip()
        if candidate and candidate.lower() not in {'theme', 'themes', 'unknown', 'unknown theme'}:
            return candidate

    theme_id = theme.get('theme_id')
    if theme_id:
        return f"{default} {theme_id}"
    return default


def _get_sentiment_analysis(year=None):
    """Get sentiment trend analysis."""
    try:
        sentiment = analyze_sentiment_trends(year=year)
        if not sentiment:
            return None
        
        # Calculate aggregate stats
        total_positive = sum(s.get('positive', 0) for s in sentiment)
        total_neutral = sum(s.get('neutral', 0) for s in sentiment)
        total_negative = sum(s.get('negative', 0) for s in sentiment)
        total = total_positive + total_neutral + total_negative
        
        if total == 0:
            return None
        
        return {
            'positive_count': total_positive,
            'neutral_count': total_neutral,
            'negative_count': total_negative,
            'total': total,
            'positive_rate': round((total_positive / total) * 100, 1),
            'neutral_rate': round((total_neutral / total) * 100, 1),
            'negative_rate': round((total_negative / total) * 100, 1),
            'trend': 'improving' if total_positive > total_negative else 'declining' if total_negative > total_positive else 'stable'
        }
    except Exception as e:
        logger.warning(f"Could not analyze sentiment: {e}")
        return None


def _get_anomalies(year=None, limit=5):
    """Fetch critical anomalies for the period."""
    try:
        rows = execute_safe_query(
            '''
            SELECT theme, month, mention_count, expected_baseline, z_score, severity
            FROM ThemeAnomalies
            WHERE COALESCE(severity, 'low') IN ('critical', 'high')
            ORDER BY z_score DESC
            LIMIT ?
            ''',
            (limit,)
        )
        return rows or []
    except Exception as e:
        logger.warning(f"Could not fetch anomalies: {e}")
        return []


def _generate_ai_recommendations(themes, sentiment_data, anomalies):
    """Generate high-quality, data-driven recommendations using LLM."""
    openai_key = os.environ.get("OPENAI_API_KEY")
    if not openai_key:
        return None

    try:
        from openai import OpenAI
        client = OpenAI(api_key=openai_key)
        
        # Prepare context for the AI
        theme_context = ", ".join([_display_theme_name(t) for t in themes[:5]])
        sentiment_summary = f"{sentiment_data['positive_rate']}% positive, {sentiment_data['negative_rate']}% negative" if sentiment_data else "Unknown"
        anomaly_context = ", ".join([f"{a.get('theme')} ({a.get('severity')})" for a in anomalies[:3]]) if anomalies else "None detected"
        
        system_prompt = """
        You are a senior governance consultant. Analyze the provided meeting data and provide 3-4 professional, actionable executive recommendations.
        Be specific, forward-looking, and grounded in the data provided.
        Format each recommendation as a concise bullet point.
        """
        
        user_prompt = f"""
        Data Snapshot:
        - Top Themes: {theme_context}
        - Sentiment: {sentiment_summary}
        - Recent Anomalies: {anomaly_context}
        
        Provide 3-4 strategic recommendations for the executive board based on this data.
        """
        
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.4,
            max_tokens=500
        )
        
        output = response.choices[0].message.content.strip()
        # Clean up bullet points if AI returned them
        recs = [r.strip('•-* ').strip() for r in output.split('\n') if r.strip()]
        return recs[:4]
    except Exception as e:
        logging.warning(f"AI recommendation generation failed: {e}")
        return None


def _generate_recommendations(themes, sentiment_data, anomalies):
    """Generate recommendations based on data (AI-first with rule-based fallback)."""
    # Try AI first for reliable, dynamic insights
    ai_recs = _generate_ai_recommendations(themes, sentiment_data, anomalies)
    if ai_recs:
        return ai_recs

    # Rule-based fallback (Static but reliable)
    recommendations = []
    
    # Sentiment-based recommendation
    if sentiment_data:
        if sentiment_data['positive_rate'] > 75:
            recommendations.append("Sentiment is highly positive. Continue current initiatives and consider scaling successful practices.")
        elif sentiment_data['positive_rate'] < 40:
            recommendations.append("Sentiment is low. Conduct stakeholder interviews to identify pain points and improvement areas.")
        else:
            recommendations.append(f"Sentiment is {sentiment_data['trend']}. Monitor trends closely and adapt engagement strategies.")
    
    # Theme-based recommendation
    if themes:
        top_theme = themes[0] if themes else None
        if top_theme:
            theme_name = _display_theme_name(top_theme)
            mention_count = top_theme.get('total_mentions', 0)
            recommendations.append(
                f"'{theme_name}' is the dominant theme with {mention_count} mentions. "
                f"Assign an owner, track follow-up actions, and review progress in the next cycle."
            )
    
    # Anomaly-based recommendation
    if anomalies:
        critical_count = len([a for a in anomalies if str(a.get('severity', '')).lower() == 'critical'])
        if critical_count > 0:
            recommendations.append(f"Alert: {critical_count} critical anomalies detected. Review flagged themes immediately before using for decision-making.")
    
    return recommendations


@report_bp.route('/api/reports/executive-summary', methods=['GET'])
@jwt_required()
def get_executive_summary():
    """
    Generate an executive summary report.
    Query params:
      - year: 2024, 2025, etc. (defaults to current year)
      - days: number of days to include (optional, defaults to all-time)
      - limit: max themes to include (default 5)
    """
    try:
        year = request.args.get('year', type=int, default=datetime.now().year)
        limit = request.args.get('limit', type=int, default=5)
        payload = _build_summary_payload(year=year, limit=limit)
        return jsonify(payload), 200
    
    except Exception as e:
        logger.error(f"Executive summary generation error: {e}")
        return jsonify({'error': 'Failed to generate summary'}), 500


@report_bp.route('/api/reports/formatted-html', methods=['GET'])
@jwt_required()
def get_formatted_html_report():
    """
    Generate a formatted HTML report for email/printing.
    Query params: year (default current year)
    """
    try:
        year = request.args.get('year', type=int, default=datetime.now().year)
        data = _build_summary_payload(year=year, limit=5)
        
        # Build HTML
        html_content = f"""
        <html>
        <head>
            <meta charset="UTF-8">
            <style>
                body {{ font-family: Arial, sans-serif; color: #333; line-height: 1.6; max-width: 900px; margin: 0 auto; padding: 20px; }}
                .header {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 30px; border-radius: 8px; margin-bottom: 30px; }}
                .header h1 {{ margin: 0; font-size: 28px; }}
                .header p {{ margin: 5px 0 0 0; opacity: 0.9; }}
                .section {{ margin-bottom: 30px; page-break-inside: avoid; }}
                .section h2 {{ color: #667eea; border-bottom: 2px solid #667eea; padding-bottom: 10px; margin-bottom: 15px; }}
                .stats-grid {{ display: grid; grid-template-columns: repeat(2, 1fr); gap: 15px; margin-bottom: 15px; }}
                .stat-card {{ background: #f5f5f5; padding: 15px; border-radius: 8px; border-left: 4px solid #667eea; }}
                .stat-label {{ font-size: 12px; color: #666; text-transform: uppercase; }}
                .stat-value {{ font-size: 24px; font-weight: bold; color: #333; margin-top: 5px; }}
                .theme-list {{ list-style: none; padding: 0; }}
                .theme-item {{ padding: 12px; border-left: 3px solid #667eea; margin-bottom: 10px; background: #f9f9f9; }}
                .theme-name {{ font-weight: bold; color: #333; }}
                .theme-detail {{ font-size: 12px; color: #666; margin-top: 5px; }}
                .sentiment-bar {{ display: flex; height: 20px; border-radius: 4px; overflow: hidden; margin: 10px 0; }}
                .sentiment-positive {{ background: #4caf50; flex: var(--positive); }}
                .sentiment-neutral {{ background: #ff9800; flex: var(--neutral); }}
                .sentiment-negative {{ background: #f44336; flex: var(--negative); }}
                .recommendation {{ padding: 12px; margin: 10px 0; background: #e3f2fd; border-left: 4px solid #2196f3; border-radius: 4px; }}
                .anomaly-alert {{ padding: 12px; margin: 10px 0; background: #fff3e0; border-left: 4px solid #ff9800; border-radius: 4px; }}
                .footer {{ color: #666; font-size: 12px; text-align: center; margin-top: 40px; border-top: 1px solid #ddd; padding-top: 20px; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>Governance Summary Report</h1>
                <p>Year: {year} | Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
            </div>
            
            <div class="section">
                <h2>Executive Summary</h2>
                <p>{data.get('executive_summary', 'No summary available.')}</p>
            </div>
            
            <div class="section">
                <h2>Key Statistics</h2>
                <div class="stats-grid">
                    <div class="stat-card">
                        <div class="stat-label">Total Meetings</div>
                        <div class="stat-value">{data.get('statistics', {}).get('total_meetings', 0)}</div>
                    </div>
                    <div class="stat-card">
                        <div class="stat-label">Unique Themes</div>
                        <div class="stat-value">{data.get('statistics', {}).get('total_themes', 0)}</div>
                    </div>
                    <div class="stat-card">
                        <div class="stat-label">Sentiment Trend</div>
                        <div class="stat-value" style="text-transform: capitalize;">{data.get('statistics', {}).get('sentiment_trend', 'Unknown')}</div>
                    </div>
                    <div class="stat-card">
                        <div class="stat-label">Critical Anomalies</div>
                        <div class="stat-value">{data.get('statistics', {}).get('critical_anomalies', 0)}</div>
                    </div>
                </div>
            </div>
            
            <div class="section">
                <h2>Top Governance Themes</h2>
                <ul class="theme-list">
        """
        
        for theme in data.get('top_themes', []):
            html_content += f"""
                    <li class="theme-item">
                        <div class="theme-name">{theme.get('theme', theme.get('name', 'Unknown'))}</div>
                        <div class="theme-detail">{theme.get('mentions', 0)} mentions • {theme.get('percentage', 0)}% of total</div>
                    </li>
            """
        
        sentiment = data.get('sentiment', {})
        if sentiment:
            html_content += f"""
            </ul>
            </div>
            
            <div class="section">
                <h2>Sentiment Analysis</h2>
                <div style="--positive: {sentiment.get('positive_rate', 0)}%; --neutral: {sentiment.get('neutral_rate', 0)}%; --negative: {sentiment.get('negative_rate', 0)}%;" class="sentiment-bar">
                    <div class="sentiment-positive"></div>
                    <div class="sentiment-neutral"></div>
                    <div class="sentiment-negative"></div>
                </div>
                <p>
                    <strong>Positive:</strong> {sentiment.get('positive_rate', 0)}% ({sentiment.get('positive_count', 0)} mentions)<br>
                    <strong>Neutral:</strong> {sentiment.get('neutral_rate', 0)}% ({sentiment.get('neutral_count', 0)} mentions)<br>
                    <strong>Negative:</strong> {sentiment.get('negative_rate', 0)}% ({sentiment.get('negative_count', 0)} mentions)
                </p>
            </div>
            """
        
        if data.get('critical_anomalies'):
            html_content += """
            <div class="section">
                <h2>Critical Anomalies</h2>
            """
            for anomaly in data.get('critical_anomalies', []):
                html_content += f"""
                    <div class="anomaly-alert">
                        <strong>{anomaly.get('theme', anomaly.get('name', 'Unknown'))}</strong> ({anomaly.get('month', 'Unknown')})<br>
                        Mentions: {anomaly.get('mentions', 0)} (expected: {anomaly.get('baseline', 0)}) | Z-score: {anomaly.get('z_score', 0)}
                    </div>
                """
            html_content += """
            </div>
            """
        
        if data.get('recommendations'):
            html_content += """
            <div class="section">
                <h2>AI-Generated Recommendations</h2>
            """
            for rec in data.get('recommendations', []):
                html_content += f'<div class="recommendation">{rec}</div>'
            html_content += """
            </div>
            """
        
        html_content += """
            <div class="footer">
                <p>This report was automatically generated by the ITDS system.</p>
            </div>
        </body>
        </html>
        """
        
        return jsonify({
            'html': html_content,
            'year': year,
            'generated_at': datetime.now(timezone.utc).isoformat()
        }), 200
    
    except Exception as e:
        logger.error(f"HTML report generation error: {e}")
        return jsonify({'error': 'Failed to generate HTML report'}), 500


@report_bp.route('/api/reports/presentation/templates', methods=['GET'])
@jwt_required()
def get_presentation_templates():
    templates = [
        {'id': key, 'label': value.get('label', key.title())}
        for key, value in PRESENTATION_THEME_STYLES.items()
    ]
    return jsonify({
        'templates': templates,
        'slide_modes': [
            {'id': 'auto', 'label': 'Auto (dynamic)'},
            {'id': 'fixed', 'label': 'Fixed (always include all core slides)'},
        ]
    }), 200


@report_bp.route('/api/reports/presentation', methods=['POST'])
@jwt_required()
def export_presentation_report():
    try:
        payload = request.get_json(silent=True) or {}
        year = _safe_int(payload.get('year'), datetime.now().year)
        limit = max(3, min(10, _safe_int(payload.get('top_n_themes'), 6)))
        include_anomalies = _as_bool(payload.get('include_anomalies'), True)
        include_speaker_notes = _as_bool(payload.get('include_speaker_notes'), True)
        include_appendix = _as_bool(payload.get('include_appendix'), False)
        slide_mode = str(payload.get('slide_mode') or 'auto').strip().lower()
        template_theme = str(payload.get('template_theme') or 'corporate').strip().lower()
        title = str(payload.get('title') or '').strip() or None

        data = _build_summary_payload(year=year, limit=limit)
        pptx_stream = _build_presentation_bytes(
            data=data,
            template_theme=template_theme,
            include_anomalies=include_anomalies,
            include_speaker_notes=include_speaker_notes,
            title=title,
            slide_mode=slide_mode,
            include_appendix=include_appendix,
        )
        safe_theme = template_theme if template_theme in PRESENTATION_THEME_STYLES else 'corporate'
        filename = f"governance-presentation-{year}-{safe_theme}.pptx"

        return send_file(
            pptx_stream,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        logger.error(f"PPTX export generation error: {e}")
        return jsonify({'error': 'Failed to generate PPTX presentation'}), 500
