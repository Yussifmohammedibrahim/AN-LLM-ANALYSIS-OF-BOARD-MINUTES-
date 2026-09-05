from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

OUT_DIR = os.path.join('uploads', 'presentations')
os.makedirs(OUT_DIR, exist_ok=True)
OUT_PATH = os.path.join(OUT_DIR, 'sample_presentation_template.pptx')

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Title slide
s = prs.slides.add_slide(prs.slide_layouts[0])
title = s.shapes.title
subtitle = s.placeholders[1]

title.text = "Executive Summary: Monthly Themes & Trends"
subtitle.text = "Auto-generated template — replace placeholders with live data"

# Agenda slide
s = prs.slides.add_slide(prs.slide_layouts[1])
shapes = s.shapes
shapes.title.text = "Agenda"
body = shapes.placeholders[1].text_frame
body.text = "1. Executive summary"

p = body.add_paragraph()
p.text = "2. Top themes and counts"
p.level = 1
p = body.add_paragraph()
p.text = "3. Trend charts"
p.level = 1
p = body.add_paragraph()
p.text = "4. Anomalies & recommendations"
p.level = 1

# Executive summary slide
s = prs.slides.add_slide(prs.slide_layouts[5])
shapes = s.shapes
shapes.title.text = "Executive Summary"
left = Inches(0.5)
top = Inches(1.3)
width = Inches(6.5)
height = Inches(4.8)

txBox = s.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Top takeaways"
p.font.size = Pt(24)

p = tf.add_paragraph()
p.text = "• Top theme: [THEME_NAME] — [TOTAL_MENTIONS] mentions"
p.level = 1
p.font.size = Pt(18)

p = tf.add_paragraph()
p.text = "• Growing topics: [TOP_GROWING_THEME] (↑ [GROWTH_RATE]%)"
p.level = 1
p.font.size = Pt(18)

# Placeholder for charts slide
s = prs.slides.add_slide(prs.slide_layouts[5])
shapes = s.shapes
shapes.title.text = "Top Themes — Chart"

left = Inches(7.1)
top = Inches(1.3)
width = Inches(5.5)
height = Inches(4.8)
chart_box = s.shapes.add_textbox(left, top, width, height)
ct = chart_box.text_frame
ct.text = "[CHART PLACEHOLDER]\nReplace with generated chart image (png)"
ct.paragraphs[0].font.size = Pt(16)
ct.paragraphs[0].alignment = PP_ALIGN.CENTER

# Anomalies & recommendations slide
s = prs.slides.add_slide(prs.slide_layouts[1])
shapes = s.shapes
shapes.title.text = "Anomalies & Recommendations"
body = shapes.placeholders[1].text_frame
body.text = "Detected anomalies"

p = body.add_paragraph()
p.text = "• [ANOMALY_1_DESCRIPTION] — suggested action: [ACTION]"
p.level = 1

p = body.add_paragraph()
p.text = "• [ANOMALY_2_DESCRIPTION] — suggested action: [ACTION]"
p.level = 1

# Closing slide
s = prs.slides.add_slide(prs.slide_layouts[6])
shapes = s.shapes
title_box = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12.0), Inches(0.8))
title_tf = title_box.text_frame
title_tf.text = "Notes & Next Steps"
title_tf.paragraphs[0].font.size = Pt(36)

prs.save(OUT_PATH)
print(f"Sample PPTX template created at: {OUT_PATH}")
