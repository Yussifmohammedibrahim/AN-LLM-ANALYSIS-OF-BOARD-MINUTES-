import unittest
from io import BytesIO

from pptx import Presentation

from itds_env.app.report_generator import (
    PRESENTATION_THEME_STYLES,
    _build_presentation_bytes,
    _build_summary_payload,
)


class PresentationExportTests(unittest.TestCase):
    def test_can_build_summary_payload(self):
        payload = _build_summary_payload(year=2026, limit=5)
        self.assertIn('year', payload)
        self.assertIn('top_themes', payload)
        self.assertIn('statistics', payload)

    def test_can_generate_all_template_themes(self):
        payload = _build_summary_payload(year=2026, limit=5)
        for theme_name in PRESENTATION_THEME_STYLES.keys():
            stream = _build_presentation_bytes(
                data=payload,
                template_theme=theme_name,
                include_anomalies=True,
                include_speaker_notes=True,
                title='Unit Test Deck',
            )
            content = stream.getvalue()
            self.assertGreater(len(content), 1024, f'PPTX too small for theme {theme_name}')

    def test_auto_mode_uses_dynamic_slide_count(self):
        sparse_payload = {
            'year': 2026,
            'executive_summary': 'Summary available, but no themes/sentiment/anomalies.',
            'statistics': {
                'total_meetings': 0,
                'total_themes': 0,
                'sentiment_trend': 'unknown',
                'critical_anomalies': 0,
            },
            'top_themes': [],
            'sentiment': {
                'positive_rate': 0,
                'neutral_rate': 0,
                'negative_rate': 0,
                'trend': 'unknown',
            },
            'critical_anomalies': [],
            'recommendations': [],
        }

        auto_stream = _build_presentation_bytes(
            data=sparse_payload,
            template_theme='corporate',
            include_anomalies=True,
            include_speaker_notes=False,
            slide_mode='auto',
            include_appendix=False,
            title='Auto Mode Deck',
        )
        fixed_stream = _build_presentation_bytes(
            data=sparse_payload,
            template_theme='corporate',
            include_anomalies=True,
            include_speaker_notes=False,
            slide_mode='fixed',
            include_appendix=False,
            title='Fixed Mode Deck',
        )

        auto_prs = Presentation(BytesIO(auto_stream.getvalue()))
        fixed_prs = Presentation(BytesIO(fixed_stream.getvalue()))
        self.assertEqual(len(auto_prs.slides), 3)
        self.assertEqual(len(fixed_prs.slides), 6)


if __name__ == '__main__':
    unittest.main()
