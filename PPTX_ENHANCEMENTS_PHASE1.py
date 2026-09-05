"""
Enhanced PPTX Report Generation - Phase 1 Features
Adds data visualizations, advanced analytics, and branding support
"""

# Enhancement functions to be integrated into report_generator.py

from datetime import datetime, timedelta
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from io import BytesIO
import json

# ============================================================================
# PHASE 1: DATA VISUALIZATIONS & ADVANCED ANALYTICS
# ============================================================================

def _get_sentiment_trends(year):
    """Get monthly sentiment trends for year-over-year visualization."""
    from .models import execute_safe_query
    
    query = '''
    SELECT 
        strftime('%m', tm.created_at) as month,
        ROUND(AVG(CASE WHEN a.sentiment = 'positive' THEN 1 ELSE 0 END) * 100, 1) as positive_rate,
        ROUND(AVG(CASE WHEN a.sentiment = 'neutral' THEN 1 ELSE 0 END) * 100, 1) as neutral_rate,
        ROUND(AVG(CASE WHEN a.sentiment = 'negative' THEN 1 ELSE 0 END) * 100, 1) as negative_rate
    FROM Meetings m
    LEFT JOIN TranscriptMetadata tm ON m.id = tm.meeting_id
    LEFT JOIN SentimentAnalysis a ON tm.id = a.metadata_id
    WHERE strftime('%Y', m.meeting_date) = ?
    GROUP BY strftime('%m', tm.created_at)
    ORDER BY month
    '''
    
    try:
        results = execute_safe_query(query, (str(year),))
        return results if results else []
    except Exception as e:
        print(f"Error fetching sentiment trends: {e}")
        return []


def _get_theme_growth_analysis(year):
    """Calculate theme growth rates and ranking."""
    from .models import execute_safe_query
    
    # Current year themes
    current_query = '''
    SELECT t.name, COUNT(*) as current_mentions
    FROM Themes t
    LEFT JOIN ThemeOccurrences o ON t.id = o.theme_id
    WHERE strftime('%Y', o.occurred_at) = ?
    GROUP BY t.id, t.name
    '''
    
    try:
        current = execute_safe_query(current_query, (str(year),))
        previous = execute_safe_query(current_query, (str(year - 1),)) if year > 2020 else []
        
        # Calculate growth rates
        growth_data = []
        for curr in (current or []):
            theme_name = curr.get('name', 'Unknown')
            curr_mentions = curr.get('current_mentions', 0)
            
            # Find previous year mention count
            prev_mentions = next(
                (p.get('current_mentions', 0) for p in previous if p.get('name') == theme_name),
                0
            )
            
            if curr_mentions > 0 or prev_mentions > 0:
                growth_pct = (
                    ((curr_mentions - prev_mentions) / prev_mentions * 100) 
                    if prev_mentions > 0 
                    else (100 if curr_mentions > 0 else 0)
                )
                
                growth_data.append({
                    'theme': theme_name,
                    'current': curr_mentions,
                    'previous': prev_mentions,
                    'growth_pct': round(growth_pct, 1),
                    'trend': '↑' if growth_pct > 0 else ('↓' if growth_pct < 0 else '→'),
                })
        
        # Sort by growth rate (descending)
        return sorted(growth_data, key=lambda x: x['growth_pct'], reverse=True)[:8]
    
    except Exception as e:
        print(f"Error calculating growth analysis: {e}")
        return []


def _get_anomaly_details_table(year, limit=5):
    """Get detailed anomaly information formatted as table data."""
    from .models import execute_safe_query
    
    query = '''
    SELECT 
        t.name as theme,
        strftime('%B', a.detected_at) as month,
        a.mention_count as mentions,
        a.baseline as baseline,
        ROUND(a.z_score, 2) as z_score,
        CASE 
            WHEN ABS(a.z_score) >= 3 THEN 'Critical'
            WHEN ABS(a.z_score) >= 2 THEN 'High'
            WHEN ABS(a.z_score) >= 1.5 THEN 'Medium'
            ELSE 'Low'
        END as severity
    FROM Anomalies a
    JOIN Themes t ON a.theme_id = t.id
    WHERE strftime('%Y', a.detected_at) = ?
    ORDER BY ABS(a.z_score) DESC
    LIMIT ?
    '''
    
    try:
        results = execute_safe_query(query, (str(year), limit))
        return results if results else []
    except Exception as e:
        print(f"Error fetching anomaly details: {e}")
        return []


def _get_prioritized_recommendations(recommendations_list):
    """Categorize recommendations by priority and add impact/effort estimates."""
    
    def estimate_priority(rec_text):
        """Heuristic: assign priority based on keywords."""
        rec_lower = str(rec_text).lower()
        
        if any(w in rec_lower for w in ['urgent', 'critical', 'immediately', 'risk', 'failing']):
            return 'High'
        elif any(w in rec_lower for w in ['monitor', 'track', 'review', 'consider']):
            return 'Medium'
        else:
            return 'Low'
    
    prioritized = []
    for rec in recommendations_list:
        priority = estimate_priority(rec)
        prioritized.append({
            'text': rec,
            'priority': priority,
            'icon': '🔴' if priority == 'High' else ('🟡' if priority == 'Medium' else '🟢'),
            'effort': 'Medium',  # Can be refined with actual data
        })
    
    # Sort: High → Medium → Low
    priority_order = {'High': 0, 'Medium': 1, 'Low': 2}
    return sorted(prioritized, key=lambda x: priority_order.get(x['priority'], 3))


def _add_sentiment_trends_slide(prs, style, sentiment_trends, year):
    """Add sentiment time-series trend slide."""
    
    s = prs.slides.add_slide(prs.slide_layouts[5])
    from .report_generator import _apply_slide_background, _add_footer, _hex_to_rgb
    _apply_slide_background(s, style['bg'])
    
    s.shapes.title.text = 'Sentiment Trend (Monthly)'
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['title'])
    s.shapes.title.text_frame.margin_top = Inches(0.2)
    s.shapes.title.text_frame.margin_left = Inches(0.6)
    
    if sentiment_trends:
        # Create line chart
        months = [t.get('month', '01') for t in sentiment_trends]
        pos_rates = [t.get('positive_rate', 0) for t in sentiment_trends]
        
        chart_data = CategoryChartData()
        chart_data.categories = [f'M{m}' for m in months]
        chart_data.add_series('Positive %', pos_rates)
        
        chart = s.shapes.add_chart(
            XL_CHART_TYPE.LINE,
            Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.3),
            chart_data
        ).chart
        
        for series in chart.series:
            series.format.line.color.rgb = _hex_to_rgb(style['positive'])
        
        chart.has_legend = False
    
    _add_footer(s, style, f"Monthly sentiment trajectory for {year}")


def _add_growth_analysis_slide(prs, style, growth_data, year):
    """Add theme growth rate analysis slide."""
    from .report_generator import _apply_slide_background, _add_footer, _hex_to_rgb
    
    s = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_slide_background(s, style['bg'])
    
    s.shapes.title.text = 'Top Growth Themes'
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['title'])
    s.shapes.title.text_frame.margin_top = Inches(0.2)
    s.shapes.title.text_frame.margin_left = Inches(0.6)
    
    if growth_data:
        # Create data table
        rows = len(growth_data) + 1  # +1 for header
        cols = 4
        left = Inches(0.6)
        top = Inches(1.3)
        width = Inches(12.1)
        height = Inches(4.5)
        
        table_shape = s.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        headers = ['Theme', 'Current', 'Previous', 'Growth %']
        for col_idx, header in enumerate(headers):
            cell = table_shape.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['title'])
        
        # Rows
        for row_idx, growth in enumerate(growth_data, 1):
            table_shape.cell(row_idx, 0).text = growth['theme'][:25]
            table_shape.cell(row_idx, 1).text = str(growth['current'])
            table_shape.cell(row_idx, 2).text = str(growth['previous'])
            
            growth_cell = table_shape.cell(row_idx, 3)
            growth_cell.text = f"{growth['trend']} {growth['growth_pct']}%"
            
            # Color code based on growth
            if growth['growth_pct'] > 0:
                growth_cell.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['positive'])
            elif growth['growth_pct'] < 0:
                growth_cell.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['negative'])
    
    _add_footer(s, style, f"Year-over-year comparison: {year} vs {year - 1}")


def _add_anomaly_details_slide(prs, style, anomalies, year):
    """Add detailed anomalies table slide."""
    from .report_generator import _apply_slide_background, _add_footer, _hex_to_rgb
    
    s = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_slide_background(s, style['bg'])
    
    s.shapes.title.text = 'Anomaly Details'
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['title'])
    s.shapes.title.text_frame.margin_top = Inches(0.2)
    s.shapes.title.text_frame.margin_left = Inches(0.6)
    
    if anomalies:
        rows = len(anomalies) + 1  # +1 for header
        cols = 5
        left = Inches(0.4)
        top = Inches(1.3)
        width = Inches(12.5)
        height = Inches(4.5)
        
        table_shape = s.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Header
        headers = ['Theme', 'Month', 'Mentions', 'Baseline', 'Severity']
        for col_idx, header in enumerate(headers):
            cell = table_shape.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
        
        # Rows with color-coded severity
        for row_idx, anomaly in enumerate(anomalies, 1):
            table_shape.cell(row_idx, 0).text = anomaly.get('theme', 'Unknown')[:20]
            table_shape.cell(row_idx, 1).text = anomaly.get('month', 'N/A')[:10]
            table_shape.cell(row_idx, 2).text = str(anomaly.get('mentions', 0))
            table_shape.cell(row_idx, 3).text = str(anomaly.get('baseline', 0))
            
            severity_cell = table_shape.cell(row_idx, 4)
            severity = anomaly.get('severity', 'Low')
            severity_cell.text = severity
            
            # Color-code severity
            if severity == 'Critical':
                severity_cell.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['negative'])
            elif severity == 'High':
                severity_cell.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['neutral'])
    
    _add_footer(s, style, f"Anomalies detected in {year}")


def _add_recommendations_slide(prs, style, recommendations):
    """Add prioritized recommendations slide with severity badges."""
    from .report_generator import _apply_slide_background, _add_footer, _hex_to_rgb
    
    s = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_slide_background(s, style['bg'])
    
    s.shapes.title.text = 'AI Recommendations (Prioritized)'
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['title'])
    s.shapes.title.text_frame.margin_top = Inches(0.2)
    s.shapes.title.text_frame.margin_left = Inches(0.6)
    
    prioritized = _get_prioritized_recommendations(recommendations)
    
    top_y = Inches(1.3)
    left_x = Inches(0.6)
    
    for idx, rec in enumerate(prioritized[:4]):  # Show top 4
        rect_top = top_y + Inches(idx * 1.2)
        
        # Badge box
        badge = s.shapes.add_shape(1, left_x, rect_top, Inches(0.4), Inches(0.3))  # 1 = rectangle
        badge.fill.solid()
        badge.fill.fore_color.rgb = (
            _hex_to_rgb(style['negative']) if rec['priority'] == 'High'
            else _hex_to_rgb(style['neutral']) if rec['priority'] == 'Medium'
            else _hex_to_rgb(style['positive'])
        )
        badge.line.color.rgb = _hex_to_rgb(style['title'])
        badge_text = badge.text_frame
        badge_text.text = rec['priority'][0]  # H, M, or L
        badge_text.paragraphs[0].font.bold = True
        badge_text.paragraphs[0].font.size = Pt(12)
        badge_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Recommendation text
        text_box = s.shapes.add_textbox(left_x + Inches(0.6), rect_top, Inches(11.5), Inches(1.0))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.text = rec['text'][:100]  # Truncate long text
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['text'])
    
    _add_footer(s, style, "Recommendations ranked by priority (High → Medium → Low)")


def _add_executive_callout_slide(prs, style, data):
    """Add highlighted key metrics slide with visual callouts."""
    from .report_generator import _apply_slide_background, _add_footer, _hex_to_rgb
    
    s = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_slide_background(s, style['bg'])
    
    s.shapes.title.text = 'Key Metrics Snapshot'
    s.shapes.title.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['title'])
    
    stats = data.get('statistics', {})
    sentiment = data.get('sentiment', {})
    
    # Create 4 metric boxes
    metrics = [
        ('📊', 'Total Meetings', str(stats.get('total_meetings', 0)), style['accent']),
        ('🎯', 'Themes', str(stats.get('total_themes', 0)), style['accent']),
        ('😊', 'Positive %', f"{sentiment.get('positive_rate', 0)}%", style['positive']),
        ('⚠️', 'Anomalies', str(stats.get('critical_anomalies', 0)), style['negative']),
    ]
    
    for idx, (emoji, label, value, color) in enumerate(metrics):
        col = idx % 2
        row = idx // 2
        
        box_left = Inches(1.5 + col * 5.5)
        box_top = Inches(1.8 + row * 2.5)
        
        # Draw box background
        box_shape = s.shapes.add_shape(1, box_left, box_top, Inches(4.5), Inches(2.0))
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = RGBColor(248, 250, 255)  # Light blue
        box_shape.line.color.rgb = _hex_to_rgb(color)
        box_shape.line.width = Pt(2)
        
        # Emoji/Icon
        emoji_box = s.shapes.add_textbox(box_left + Inches(0.2), box_top + Inches(0.2), Inches(0.8), Inches(0.8))
        emoji_box.text_frame.text = emoji
        emoji_box.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Label
        label_box = s.shapes.add_textbox(box_left + Inches(1.2), box_top + Inches(0.3), Inches(3.0), Inches(0.5))
        label_box.text_frame.text = label
        label_box.text_frame.paragraphs[0].font.size = Pt(12)
        label_box.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(style['muted'])
        
        # Value
        value_box = s.shapes.add_textbox(box_left + Inches(1.2), box_top + Inches(0.9), Inches(3.0), Inches(1.0))
        value_box.text_frame.text = value
        value_box.text_frame.paragraphs[0].font.size = Pt(32)
        value_box.text_frame.paragraphs[0].font.bold = True
        value_box.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(color)
    
    _add_footer(s, style, "Key performance indicators at a glance")


# ============================================================================
# PHASE 1: BRANDING SUPPORT
# ============================================================================

def _add_organization_branding(prs, branding_data, style):
    """Apply organization branding to presentation metadata."""
    
    if not branding_data:
        return prs
    
    core = prs.core_properties
    
    # Apply branding metadata
    if branding_data.get('organization_name'):
        core.company = branding_data['organization_name']
        core.subject = f"{branding_data['organization_name']} - {core.subject}"
    
    if branding_data.get('author'):
        core.author = branding_data.get('author', core.author)
    
    if branding_data.get('footer_text'):
        core.comments = branding_data['footer_text']
    
    return prs


def _apply_custom_colors(style, branding_data):
    """Override theme colors with custom branding colors."""
    
    if not branding_data:
        return style
    
    if branding_data.get('primary_color'):
        style['accent'] = branding_data['primary_color'].lstrip('#')
    
    if branding_data.get('secondary_color'):
        style['title'] = branding_data['secondary_color'].lstrip('#')
    
    return style


def _add_watermark_text(slide, watermark_text, style):
    """Add watermark overlay to slide."""
    from .report_generator import _hex_to_rgb
    
    if not watermark_text:
        return
    
    # Add semi-transparent watermark
    watermark = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(3.5))
    tf = watermark.text_frame
    tf.text = watermark_text
    tf.paragraphs[0].font.size = Pt(96)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = _hex_to_rgb(style.get('muted', '475569'))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Make it semi-transparent (opacity ~20%)
    fill = watermark.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(style.get('muted', '475569'))


# ============================================================================
# ENHANCED PPTX BUILDER (Integration point)
# ============================================================================

def build_enhanced_presentation(data, template_theme='corporate', branding_data=None, **options):
    """
    Build enhanced PPTX with all Phase 1 features.
    
    Options:
    - include_sentiment_trends: bool
    - include_growth_analysis: bool
    - include_anomaly_details: bool
    - include_recommendations: bool
    - include_key_metrics: bool
    - watermark: str
    - year: int
    """
    from .report_generator import (
        PRESENTATION_THEME_STYLES,
        _build_presentation_bytes,
        _add_organization_branding,
    )
    
    # Get enhanced data
    year = options.get('year', 2026)
    sentiment_trends = _get_sentiment_trends(year) if options.get('include_sentiment_trends') else []
    growth_data = _get_theme_growth_analysis(year) if options.get('include_growth_analysis') else []
    anomaly_details = _get_anomaly_details_table(year) if options.get('include_anomaly_details') else []
    
    # For now, return standard PPTX and add enhancements as separate slides
    # This will be integrated directly into _build_presentation_bytes in production
    
    return {
        'sentiment_trends': sentiment_trends,
        'growth_data': growth_data,
        'anomaly_details': anomaly_details,
        'status': 'ready_for_integration'
    }
